from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import os
def _create_modifier_from_attributes(ref_attribute, targ_attribute):
    return ref_attribute + " and " + targ_attribute
def _cat_captions(caps, dataset):
    I = []
    if dataset == 'fashionIQ':
        for i in range(len(caps)):
            if i % 2 == 0:
                I.append(_create_modifier_from_attributes(caps[i], caps[i+1]))
            else:
                I.append(_create_modifier_from_attributes(caps[i], caps[i-1]))
    elif dataset == 'shoes':
        for i in range(len(caps)):
            I.append(caps[i])
    return I
def caption_preprocessing(caption_file, dataset):
    caps = [ff.strip('\n').split(';')[-1] for ff in caption_file]
    caps_cat = _cat_captions(caps, dataset)

    return caps_cat

if __name__ == '__main__':

    dataset = 'shoes'
    experiment_name = 'shoes_CosMo_eval'

    if dataset == 'fashionIQ':
        domains = ['shirt']

    elif dataset == 'shoes':
        domains = ['shoes']

    for DOMAIN in domains:
        print(f"Extract the images of {DOMAIN} domain")

        parent_path = f'/home/jaehyun98/git/uncertain_retrieval/visualize_result/{experiment_name}/{DOMAIN}/wrong/'

        if dataset == 'fashionIQ':
            proper_caption_path = f'/home/jaehyun98/git/uncertain_retrieval/data/{dataset}/captions_pairs/fashion_iq-val-cap-{DOMAIN}.txt'
            ambiguous_caption_path = f'/home/jaehyun98/git/uncertain_retrieval/data/{dataset}/captions_pairs/ambiguous_fashion_iq-val-cap-{DOMAIN}.txt'

            with open(proper_caption_path, 'r') as f:
                proper_caption_file = f.readlines()

            cat_proper_caption = caption_preprocessing(proper_caption_file, dataset)

            with open(ambiguous_caption_path, 'r') as f:
                ambiguous_caption_file = f.readlines()

            cat_ambiguous_caption = caption_preprocessing(ambiguous_caption_file, dataset)

            print(len(cat_proper_caption), len(cat_ambiguous_caption))

        elif dataset == 'shoes':
            proper_caption_path = f'/home/jaehyun98/git/uncertain_retrieval/data/{dataset}/shoes-cap-test.txt'
            # ambiguous_caption_path = f'/home/jaehyun98/git/uncertain_retrieval/data/{dataset}/ambiguous_fashion_iq-val-cap-{DOMAIN}.txt'

            with open(proper_caption_path, 'r') as f:
                proper_caption_file = f.readlines()

            cat_proper_caption = caption_preprocessing(proper_caption_file, dataset)

        all_indexes = os.listdir(parent_path)

        print(f'number of indexes : {len(all_indexes)}')

        flag = 0
        prs = Presentation()
        for cur_index in all_indexes:
            # print(f'cur_index : {cur_index}')
            cur_path = os.path.join(parent_path, cur_index)

            source_image_path = os.path.join(cur_path, 'source_image')
            answer_image_path = os.path.join(cur_path, 'answer_image')
            # retrieved_proper_image_path = os.path.join(cur_path, 'retrieved_proper_images')
            retrieved_image_path = os.path.join(cur_path, 'wrong_retrieved_images')

            source_image = os.listdir(source_image_path)[0]
            answer_image = os.listdir(answer_image_path)[0]
            # retrieved_proper_images = os.listdir(retrieved_proper_image_path)
            retrieved_images = os.listdir(retrieved_image_path)

            # Source Image
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Choosing a blank slide layout
            left = Inches(0.1)  # 1 inch from the left of the slide
            top = Inches(0.1)  # 1 inch from the top of the slide
            height = Inches(1)  # Height of the image
            pic = slide.shapes.add_picture(os.path.join(source_image_path, source_image), left, top, height=height)

            # Target Image
            left = Inches(2.1)  # 1 inch from the left of the slide
            top = Inches(0.1)  # 1 inch from the top of the slide
            height = Inches(1)  # Height of the image
            pic = slide.shapes.add_picture(os.path.join(answer_image_path, answer_image), left, top, height=height)

            # Add source proper text
            tb1_left = Inches(4.1)
            tb1_top = Inches(0.1)
            tb1_width = Inches(5)
            tb1_height = Inches(1)
            textbox1 = slide.shapes.add_textbox(tb1_left, tb1_top, tb1_width, tb1_height)
            tf1 = textbox1.text_frame
            tf1.text = cat_proper_caption[int(cur_index)]

            # # Add source ambiguous text
            # tb1_left = Inches(4.1)
            # tb1_top = Inches(0.8)
            # tb1_width = Inches(5)
            # tb1_height = Inches(1)
            # textbox1 = slide.shapes.add_textbox(tb1_left, tb1_top, tb1_width, tb1_height)
            # tf1 = textbox1.text_frame
            # tf1.text = cat_ambiguous_caption[int(cur_index)]

            # # Add text
            # tb1_left = Inches(0.1)
            # tb1_top = Inches(2.5)
            # tb1_width = Inches(5)
            # tb1_height = Inches(1)
            # textbox1 = slide.shapes.add_textbox(tb1_left, tb1_top, tb1_width, tb1_height)
            # tf1 = textbox1.text_frame
            # tf1.text = "Retrieved Images with Proper Text"
            #
            # # Retrieved Images
            # for idx, cur_retrieved_proper_image in enumerate(retrieved_proper_images):
            #     left = Inches(idx)  # 1 inch from the left of the slide
            #     top = Inches(3.1)  # 1 inch from the top of the slide
            #     height = Inches(1)  # Height of the image
            #     pic = slide.shapes.add_picture(os.path.join(retrieved_proper_image_path, cur_retrieved_proper_image), left, top, height=height)

            # Add text
            tb1_left = Inches(0.1)
            tb1_top = Inches(4.5)
            tb1_width = Inches(5)
            tb1_height = Inches(1)
            textbox1 = slide.shapes.add_textbox(tb1_left, tb1_top, tb1_width, tb1_height)
            tf1 = textbox1.text_frame
            tf1.text = "Retrieved Images"

            # Retrieved Images
            for idx, cur_retrieved_image in enumerate(retrieved_images):
                left = Inches(idx)  # 1 inch from the left of the slide
                top = Inches(5.1)  # 1 inch from the top of the slide
                height = Inches(1)  # Height of the image
                pic = slide.shapes.add_picture(os.path.join(retrieved_image_path, cur_retrieved_image), left, top, height=height)

            # Add text for index of such data
            tb1_left = Inches(0.1)
            tb1_top = Inches(6)
            tb1_width = Inches(5)
            tb1_height = Inches(1)
            textbox1 = slide.shapes.add_textbox(tb1_left, tb1_top, tb1_width, tb1_height)
            tf1 = textbox1.text_frame
            tf1.text = f"Index of data : {cur_index}"

            flag += 1

            if flag == 100:
                break

        # Save the presentation
        prs.save(f'{experiment_name}_wrong_{dataset}_{DOMAIN}_analysis.pptx')